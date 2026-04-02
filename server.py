from __future__ import annotations

import json
import os
import queue
import socket
import subprocess
import threading
import time
import uuid
import sys
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Optional

from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.responses import FileResponse, JSONResponse, StreamingResponse, Response
import mimetypes
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, Field


BASE_DIR = Path(__file__).resolve().parent
WEB_DIR = BASE_DIR / "web"
OUTPUTS_DIR = BASE_DIR / "outputs"
OUTPUTS_DIR.mkdir(parents=True, exist_ok=True)
UPLOADS_DIR = BASE_DIR / "uploads"
UPLOADS_DIR.mkdir(parents=True, exist_ok=True)

def _find_best_python() -> str:
    """Find the best Python executable: prefer figure conda env, fallback to sys.executable."""
    explicit = os.environ.get("AUTOFIGURE_PYTHON")
    if explicit and Path(explicit).is_file():
        return explicit

    # Try to find figure conda env Python
    conda_exe = os.environ.get("CONDA_EXE")
    if conda_exe:
        conda_root = Path(conda_exe).parent.parent
        for env_name in ["figure"]:
            candidates = [
                conda_root / "envs" / env_name / "python.exe",
                conda_root / "envs" / env_name / "bin" / "python",
            ]
            for candidate in candidates:
                if candidate.is_file():
                    return str(candidate)

    # Also check CONDA_PREFIX siblings
    conda_prefix = os.environ.get("CONDA_PREFIX")
    if conda_prefix:
        envs_dir = Path(conda_prefix).parent
        for env_name in ["figure"]:
            for suffix in ["python.exe", "bin/python"]:
                candidate = envs_dir / env_name / suffix
                if candidate.is_file():
                    return str(candidate)

    return sys.executable


PYTHON_EXECUTABLE = _find_best_python()

DEFAULT_SAM_PROMPT = "icon, illustration, person, robot, machine, device, animal, Signal, MRI"
DEFAULT_PLACEHOLDER_MODE = "label"
DEFAULT_MERGE_THRESHOLD = 0.3
DEFAULT_MAX_BOX_AREA_RATIO = 0.1
DEFAULT_SAM_MIN_SCORE = 0.2

SVG_EDIT_CANDIDATES = [
    ("vendor/svg-edit/editor/index.html", WEB_DIR / "vendor" / "svg-edit" / "editor" / "index.html"),
    ("vendor/svg-edit/editor.html", WEB_DIR / "vendor" / "svg-edit" / "editor.html"),
    ("vendor/svg-edit/index.html", WEB_DIR / "vendor" / "svg-edit" / "index.html"),
]


def _resolve_svg_edit_path() -> tuple[bool, str | None]:
    for rel, path in SVG_EDIT_CANDIDATES:
        if path.is_file():
            return True, f"/{rel}"
    return False, None


@dataclass
class Job:
    job_id: str
    output_dir: Path
    process: subprocess.Popen
    queue: queue.Queue
    log_path: Path
    log_lock: threading.Lock = field(default_factory=threading.Lock)
    seen: set[str] = field(default_factory=set)
    done: bool = False

    def push(self, event: str, data: dict) -> None:
        self.queue.put({"event": event, "data": data})

    def write_log(self, stream: str, line: str) -> None:
        with self.log_lock:
            with open(self.log_path, "a", encoding="utf-8") as handle:
                handle.write(f"[{stream}] {line}\n")


class RunRequest(BaseModel):
    method_text: str = Field(..., min_length=1)
    image_model: Optional[str] = None
    svg_model: Optional[str] = None
    fix_svg_model: Optional[str] = None
    image_provider: str = "gemini"
    image_api_key: str = ""
    image_base_url: str = ""
    svg_provider: str = "gemini"
    svg_api_key: str = ""
    svg_base_url: str = ""
    fix_svg_provider: str = "gemini"
    fix_svg_api_key: str = ""
    fix_svg_base_url: str = ""
    sam_prompt: Optional[str] = None
    sam_backend: Optional[str] = None
    sam_api_key: Optional[str] = None
    sam_max_masks: Optional[int] = None
    sam_min_score: Optional[float] = None
    placeholder_mode: Optional[str] = None
    merge_threshold: Optional[float] = None
    max_box_area_ratio: Optional[float] = None
    optimize_iterations: Optional[int] = None
    reference_image_path: Optional[str] = None
    input_image_path: Optional[str] = None
    resume_dir: Optional[str] = None


app = FastAPI()

JOBS: dict[str, Job] = {}


@app.get("/api/config")
def get_config() -> JSONResponse:
    available, rel_path = _resolve_svg_edit_path()
    return JSONResponse({"svgEditAvailable": available, "svgEditPath": rel_path})


@app.post("/api/run")
def run_job(req: RunRequest) -> JSONResponse:
    job_id = datetime.now().strftime("%Y%m%d_%H%M%S_") + uuid.uuid4().hex[:8]
    output_dir = OUTPUTS_DIR / job_id
    output_dir.mkdir(parents=True, exist_ok=True)

    cmd = [
        PYTHON_EXECUTABLE,
        str(BASE_DIR / "autofigure2.py"),
        "--method_text",
        req.method_text,
        "--output_dir",
        str(output_dir),
    ]

    if req.image_model:
        cmd += ["--image_model", req.image_model]
    if req.svg_model:
        cmd += ["--svg_model", req.svg_model]
    if req.image_provider:
        cmd += ["--image_provider", req.image_provider]
    if req.image_api_key:
        cmd += ["--image_api_key", req.image_api_key]
    if req.image_base_url:
        cmd += ["--image_base_url", req.image_base_url]
    if req.svg_provider:
        cmd += ["--svg_provider", req.svg_provider]
    if req.svg_api_key:
        cmd += ["--svg_api_key", req.svg_api_key]
    if req.svg_base_url:
        cmd += ["--svg_base_url", req.svg_base_url]
    fix_svg_provider = req.fix_svg_provider or req.svg_provider or ""
    fix_svg_api_key = req.fix_svg_api_key or req.svg_api_key or ""
    fix_svg_base_url = req.fix_svg_base_url or req.svg_base_url or ""
    fix_svg_model = req.fix_svg_model or req.svg_model or ""
    if fix_svg_provider:
        cmd += ["--fix_svg_provider", fix_svg_provider]
    if fix_svg_api_key:
        cmd += ["--fix_svg_api_key", fix_svg_api_key]
    if fix_svg_base_url:
        cmd += ["--fix_svg_base_url", fix_svg_base_url]
    if fix_svg_model:
        cmd += ["--fix_svg_model", fix_svg_model]

    sam_prompt = req.sam_prompt or DEFAULT_SAM_PROMPT
    placeholder_mode = req.placeholder_mode or DEFAULT_PLACEHOLDER_MODE
    merge_threshold = (
        req.merge_threshold if req.merge_threshold is not None else DEFAULT_MERGE_THRESHOLD
    )
    max_box_area_ratio = (
        req.max_box_area_ratio if req.max_box_area_ratio is not None else DEFAULT_MAX_BOX_AREA_RATIO
    )
    sam_min_score = (
        req.sam_min_score if req.sam_min_score is not None else DEFAULT_SAM_MIN_SCORE
    )

    cmd += ["--sam_prompt", sam_prompt]
    cmd += ["--placeholder_mode", placeholder_mode]
    cmd += ["--merge_threshold", str(merge_threshold)]
    cmd += ["--max_box_area_ratio", str(max_box_area_ratio)]
    cmd += ["--min_score", str(sam_min_score)]
    if req.sam_backend:
        cmd += ["--sam_backend", req.sam_backend]
    if req.sam_api_key:
        cmd += ["--sam_api_key", req.sam_api_key]
    if req.sam_max_masks is not None:
        cmd += ["--sam_max_masks", str(req.sam_max_masks)]
    if req.optimize_iterations is not None:
        cmd += ["--optimize_iterations", str(req.optimize_iterations)]

    reference_path = req.reference_image_path
    if reference_path:
        reference_path = (
            str((BASE_DIR / reference_path).resolve())
            if not Path(reference_path).is_absolute()
            else reference_path
        )
        cmd += ["--reference_image_path", reference_path]

    input_image_path = req.input_image_path
    if input_image_path:
        input_image_path = (
            str((BASE_DIR / input_image_path).resolve())
            if not Path(input_image_path).is_absolute()
            else input_image_path
        )
        cmd += ["--input_image", input_image_path]

    if req.resume_dir:
        resume_path = req.resume_dir
        if not Path(resume_path).is_absolute():
            resume_path = str((OUTPUTS_DIR / resume_path).resolve())
        cmd += ["--resume_dir", resume_path]

    env = os.environ.copy()
    env["PYTHONUNBUFFERED"] = "1"

    log_path = output_dir / "run.log"
    log_path.write_text(
        f"[meta] python={PYTHON_EXECUTABLE}\n[meta] cmd={' '.join(cmd)}\n",
        encoding="utf-8",
    )

    process = subprocess.Popen(
        cmd,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        bufsize=1,
        env=env,
        cwd=str(BASE_DIR),
    )

    job = Job(
        job_id=job_id,
        output_dir=output_dir,
        process=process,
        queue=queue.Queue(),
        log_path=log_path,
    )
    JOBS[job_id] = job

    monitor_thread = threading.Thread(target=_monitor_job, args=(job,), daemon=True)
    monitor_thread.start()

    return JSONResponse({"job_id": job_id})


@app.post("/api/upload")
async def upload_reference(file: UploadFile = File(...)) -> JSONResponse:
    if not file.filename:
        raise HTTPException(status_code=400, detail="No file provided")
    if not file.content_type or not file.content_type.startswith("image/"):
        raise HTTPException(status_code=400, detail="File must be an image")

    ext = Path(file.filename).suffix.lower()
    if ext not in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif"}:
        ext = ".png"

    data = await file.read()
    if len(data) > 20 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="File too large")

    name = f"{uuid.uuid4().hex}{ext}"
    out_path = UPLOADS_DIR / name
    out_path.write_bytes(data)

    rel_path = out_path.relative_to(BASE_DIR).as_posix()
    return JSONResponse(
        {"path": rel_path, "url": f"/api/uploads/{name}", "name": file.filename}
    )


@app.get("/api/events/{job_id}")
def stream_events(job_id: str) -> StreamingResponse:
    job = JOBS.get(job_id)
    if not job:
        raise HTTPException(status_code=404, detail="Job not found")

    def event_stream():
        while True:
            try:
                item = job.queue.get(timeout=1.0)
            except queue.Empty:
                if job.done:
                    break
                continue
            if item.get("event") == "close":
                break
            yield _format_sse(item["event"], item["data"])

    return StreamingResponse(event_stream(), media_type="text/event-stream")


@app.get("/api/artifacts/{job_id}/{path:path}")
def get_artifact(job_id: str, path: str) -> Response:
    job = JOBS.get(job_id)
    if not job:
        raise HTTPException(status_code=404, detail="Job not found")

    candidate = (job.output_dir / path).resolve()
    if not str(candidate).startswith(str(job.output_dir.resolve())):
        raise HTTPException(status_code=400, detail="Invalid path")
    if not candidate.is_file():
        raise HTTPException(status_code=404, detail="File not found")

    # 为了避免运行中被持续追加的文件 (如 run.log) 
    # 被 FastAPI/Starlette 的 FileResponse 预计算出较小的 Content-Length
    # 导致读取时传输的数据大于 Content-Length 而触发 h11._util.LocalProtocolError
    # 这里直接将较小文件（<15MB）读入内存作为普通 Response 返回
    size_limit = 15 * 1024 * 1024
    if candidate.stat().st_size < size_limit:
        mime_type, _ = mimetypes.guess_type(candidate)
        return Response(content=candidate.read_bytes(), media_type=mime_type)
    
    return FileResponse(candidate)


class SaveSvgRequest(BaseModel):
    svg: str


@app.post("/api/save-svg/{job_id}")
async def save_svg(job_id: str, req: SaveSvgRequest) -> JSONResponse:
    job = JOBS.get(job_id)
    if job:
        output_dir = job.output_dir
    else:
        # Job may no longer be in memory (server restart / completed) — resolve from disk
        output_dir = OUTPUTS_DIR / job_id
        if not output_dir.is_dir():
            raise HTTPException(status_code=404, detail="Job not found")
    svg_path = output_dir / "final.svg"
    svg_path.write_text(req.svg, encoding="utf-8")
    return JSONResponse({"ok": True, "path": str(svg_path)})


class ExportPptxRequest(BaseModel):
    svg: str
    width_cm: float = 33.87
    height_cm: float = 19.05


@app.post("/api/export-pptx/{job_id}")
async def export_pptx(job_id: str, req: ExportPptxRequest) -> Response:
    """Export SVG as a PPTX with the SVG embedded as a native Office SVG image.

    Office 2016+ (and LibreOffice 6+) natively renders SVG images placed via the
    OOXML SVG extension (<asvg:svgBlip>).  The shape stays fully vector and every
    text / path element is editable inside PowerPoint without any rasterisation.

    How it works:
      1. Build a minimal PPTX in-memory with python-pptx.
      2. Add a 1×1 px placeholder PNG (required by OOXML – the blipFill must have a
         fallback raster Part even when an SVG extension is present).
      3. Inject the SVG bytes as a new Part into the zip and attach it to the shape
         via the <asvg:svgBlip r:embed="…"/> extension element.
    The resulting file opens in PPT as a fully editable SVG image.
    """
    job = JOBS.get(job_id)
    if not job:
        # Allow export even when job is no longer in memory (server restart / completed)
        output_dir = OUTPUTS_DIR / job_id
        if not output_dir.is_dir():
            raise HTTPException(status_code=404, detail="Job not found")

    try:
        import io as _io
        import zipfile as _zipfile
        from xml.etree import ElementTree as ET

        from pptx import Presentation
        from pptx.util import Cm
        from lxml import etree as _etree

        NS_A    = "http://schemas.openxmlformats.org/drawingml/2006/main"
        NS_ASVG = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
        NS_R    = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        NS_REL  = "http://schemas.openxmlformats.org/package/2006/relationships"
        NS_CT   = "http://schemas.openxmlformats.org/package/2006/content-types"
        SVG_IMG_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"

        # ── 1. Parse SVG dimensions ──────────────────────────────────────────
        svg_bytes = req.svg.encode("utf-8")
        try:
            root = ET.fromstring(svg_bytes)
            raw_w = root.get("width", "") or ""
            raw_h = root.get("height", "") or ""
            svg_w = float(raw_w.replace("px", "").strip()) if raw_w.replace("px", "").strip() else None
            svg_h = float(raw_h.replace("px", "").strip()) if raw_h.replace("px", "").strip() else None
            if svg_w is None or svg_h is None:
                vb = root.get("viewBox", "")
                parts = vb.split()
                if len(parts) == 4:
                    svg_w = float(parts[2])
                    svg_h = float(parts[3])
                else:
                    svg_w, svg_h = 1200.0, 800.0
        except Exception:
            svg_w, svg_h = 1200.0, 800.0

        # ── 2. Build PPTX with 1×1 transparent PNG as raster fallback ────────
        prs = Presentation()
        prs.slide_width = Cm(req.width_cm)
        prs.slide_height = Cm(req.height_cm)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_w_emu = int(prs.slide_width)
        slide_h_emu = int(prs.slide_height)

        TRANSPARENT_1PX_PNG = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00"
            b"\x01\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\nIDATx\x9cc"
            b"\x00\x01\x00\x00\x05\x00\x01\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        slide.shapes.add_picture(_io.BytesIO(TRANSPARENT_1PX_PNG), 0, 0, slide_w_emu, slide_h_emu)

        buf = _io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        # ── 3. Zip-level injection of SVG part + relationship + svgBlip ───────
        SLIDE_XML  = "ppt/slides/slide1.xml"
        SLIDE_RELS = "ppt/slides/_rels/slide1.xml.rels"
        SVG_PART   = "ppt/media/image_svg1.svg"
        SVG_RID    = "rIdSVG1"
        rid = SVG_RID  # pre-initialize before zip loop

        out_buf = _io.BytesIO()
        with _zipfile.ZipFile(buf, "r") as zin, \
             _zipfile.ZipFile(out_buf, "w", _zipfile.ZIP_DEFLATED) as zout:

            for item in zin.infolist():
                data = zin.read(item.filename)

                # ── patch slide rels: add SVG relationship ──────────────────
                if item.filename == SLIDE_RELS:
                    rels_tree = _etree.fromstring(data)
                    # avoid duplicate
                    existing_ids = {el.get("Id") for el in rels_tree}
                    rid = SVG_RID
                    n = 1
                    while rid in existing_ids:
                        n += 1
                        rid = f"rIdSVG{n}"
                    rel_el = _etree.SubElement(rels_tree, f"{{{NS_REL}}}Relationship")
                    rel_el.set("Id", rid)
                    rel_el.set("Type", SVG_IMG_TYPE)
                    rel_el.set("Target", "../media/image_svg1.svg")
                    data = _etree.tostring(rels_tree, xml_declaration=True,
                                           encoding="UTF-8", standalone=True)

                # ── patch slide XML: inject <asvg:svgBlip> into <a:blip> ────
                elif item.filename == SLIDE_XML:
                    slide_tree = _etree.fromstring(data)
                    # The blip is always <a:blip> regardless of parent namespace
                    for blip_el in slide_tree.iter(f"{{{NS_A}}}blip"):
                        ext_lst = blip_el.find(f"{{{NS_A}}}extLst")
                        if ext_lst is None:
                            ext_lst = _etree.SubElement(blip_el, f"{{{NS_A}}}extLst")
                        ext = _etree.SubElement(ext_lst, f"{{{NS_A}}}ext")
                        ext.set("uri", "{96DAC541-7B7A-43D3-8B79-37D633B846F1}")
                        svg_blip = _etree.SubElement(ext, f"{{{NS_ASVG}}}svgBlip")
                        svg_blip.set(f"{{{NS_R}}}embed", rid)
                        break  # only first picture
                    data = _etree.tostring(slide_tree, xml_declaration=True,
                                           encoding="UTF-8", standalone=True)

                # ── patch content types: register .svg if missing ───────────
                elif item.filename == "[Content_Types].xml":
                    ct_tree = _etree.fromstring(data)
                    if not any(el.get("Extension") == "svg"
                               for el in ct_tree.findall(f"{{{NS_CT}}}Default")):
                        d = _etree.SubElement(ct_tree, f"{{{NS_CT}}}Default")
                        d.set("Extension", "svg")
                        d.set("ContentType", "image/svg+xml")
                    data = _etree.tostring(ct_tree, xml_declaration=True,
                                           encoding="UTF-8", standalone=True)

                zout.writestr(item, data)

            # write the SVG part itself
            zout.writestr(SVG_PART, svg_bytes)

        out_buf.seek(0)
        filename = f"figure_{job_id}.pptx"
        return Response(
            content=out_buf.read(),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )
    except ImportError as exc:
        raise HTTPException(
            status_code=500,
            detail=f"python-pptx not installed. Run: pip install python-pptx lxml. ({exc})",
        )
    except Exception as exc:
        import traceback as _tb
        raise HTTPException(status_code=500, detail=f"PPTX export error: {exc}\n{_tb.format_exc()}")


@app.get("/api/uploads/{filename}")
def get_upload(filename: str) -> FileResponse:
    candidate = (UPLOADS_DIR / filename).resolve()
    if not str(candidate).startswith(str(UPLOADS_DIR.resolve())):
        raise HTTPException(status_code=400, detail="Invalid path")
    if not candidate.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(candidate)


def _format_sse(event: str, data: dict) -> str:
    payload = json.dumps(data, ensure_ascii=True)
    return f"event: {event}\ndata: {payload}\n\n"


def _monitor_job(job: Job) -> None:
    job.push("status", {"state": "started"})

    stdout_thread = threading.Thread(
        target=_pipe_output, args=(job, job.process.stdout, "stdout"), daemon=True
    )
    stderr_thread = threading.Thread(
        target=_pipe_output, args=(job, job.process.stderr, "stderr"), daemon=True
    )
    stdout_thread.start()
    stderr_thread.start()

    idle_cycles = 0
    while True:
        _scan_artifacts(job)

        if job.process.poll() is not None:
            idle_cycles += 1
        else:
            idle_cycles = 0

        if idle_cycles >= 4:
            break
        time.sleep(0.5)

    _scan_artifacts(job)
    job.push("status", {"state": "finished", "code": job.process.returncode})
    job.push(
        "artifact",
        {
            "kind": "log",
            "name": job.log_path.name,
            "path": job.log_path.relative_to(job.output_dir).as_posix(),
            "url": f"/api/artifacts/{job.job_id}/{job.log_path.name}",
        },
    )
    job.done = True
    job.push("close", {})


def _pipe_output(job: Job, pipe, stream_name: str) -> None:
    if pipe is None:
        return
    for line in iter(pipe.readline, ""):
        text = line.rstrip()
        if text:
            job.write_log(stream_name, text)
            job.push("log", {"stream": stream_name, "line": text})
    pipe.close()


def _scan_artifacts(job: Job) -> None:
    output_dir = job.output_dir
    candidates = [
        output_dir / "figure.png",
        output_dir / "samed.png",
        output_dir / "template.svg",
        output_dir / "final.svg",
    ]

    icons_dir = output_dir / "icons"
    if icons_dir.is_dir():
        candidates.extend(icons_dir.glob("icon_*.png"))

    for path in candidates:
        if not path.is_file():
            continue
        rel_path = path.relative_to(output_dir).as_posix()
        if rel_path in job.seen:
            continue
        job.seen.add(rel_path)

        kind = _classify_artifact(rel_path)
        job.push(
            "artifact",
            {
                "kind": kind,
                "name": path.name,
                "path": rel_path,
                "url": f"/api/artifacts/{job.job_id}/{rel_path}",
            },
        )


def _classify_artifact(rel_path: str) -> str:
    if rel_path == "figure.png":
        return "figure"
    if rel_path == "samed.png":
        return "samed"
    if rel_path.endswith("_nobg.png"):
        return "icon_nobg"
    if rel_path.startswith("icons/") and rel_path.endswith(".png"):
        return "icon_raw"
    if rel_path == "template.svg":
        return "template_svg"
    if rel_path == "final.svg":
        return "final_svg"
    return "artifact"


app.mount("/", StaticFiles(directory=WEB_DIR, html=True), name="static")


if __name__ == "__main__":
    import uvicorn

    def find_available_port(start_port: int, max_attempts: int = 100) -> int:
        for port in range(start_port, start_port + max_attempts):
            # Quick check: if we can connect to 127.0.0.1:port, treat it as in-use.
            with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as probe:
                probe.settimeout(0.5)
                try:
                    if probe.connect_ex(("127.0.0.1", port)) == 0:
                        print(f"Port {port} already serving on 127.0.0.1, trying next...")
                        continue
                except OSError:
                    pass

            with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as sock:
                # On Windows, SO_EXCLUSIVEADDRUSE prevents successful bind when another process
                # already owns the port (important when another app listens on 127.0.0.1).
                if os.name == "nt" and hasattr(socket, "SO_EXCLUSIVEADDRUSE"):
                    sock.setsockopt(socket.SOL_SOCKET, socket.SO_EXCLUSIVEADDRUSE, 1)
                try:
                    sock.bind(("0.0.0.0", port))
                    return port
                except OSError:
                    print(f"Port {port} is in use, trying next...")
                    continue
        raise IOError(f"No available ports found in range ({start_port} - {start_port + max_attempts})")

    def start_server_with_fallback() -> int:
        # Allow overriding the initial port from environment for flexibility in Docker/PM2/etc.
        initial_port = int(os.environ.get("AUTOFIGURE_PORT", 8000))

        try:
            actual_port = find_available_port(initial_port)
        except Exception as e:
            print(f"Startup failed while searching port: {e}")
            raise

        for attempt in range(3):
            try:
                print("--- Starting Server ---")
                print(f"Local access: http://127.0.0.1:{actual_port}")
                print("-----------------------")

                uvicorn.run(
                    "server:app",
                    host="0.0.0.0",
                    port=actual_port,
                    reload=False,
                    access_log=False,
                )
                return actual_port
            except OSError as exc:  # Catch bind errors even if a race occurs after our probe
                print(f"Failed to bind port {actual_port} ({exc}), retrying with next port...")
                actual_port += 1
        raise SystemExit("Failed to start server after multiple port attempts")

    try:
        start_server_with_fallback()
    except Exception as e:
        print(f"Startup failed: {e}")
        sys.exit(1)
